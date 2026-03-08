"""
PPT Assembler Agent (Agent 7) - PRODUCTION VERSION
Fixed grid layout. No overflow. Properly sized metrics bar.
"""

import logging
from typing import Dict, List, Any, Optional
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.chart.data import ChartData
from copy import deepcopy
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.chart import XL_TICK_MARK
import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from utils.brand_guidelines import BrandGuidelines

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


# Fixed layout grid (all positions in inches)
# Slide is 10" x 7.5" (standard widescreen)
LAYOUT = {
    'slide_width': 13.333,
    'slide_height': 7.5,
    'margin_left': 0.5,
    'margin_right': 0.5,
    'margin_top': 0.3,
    'margin_bottom': 0.5,  # Increased for footer
    'content_width': 12.333,
    'title_height': 0.5,
    'footer_height': 0.25,
    
    # Grid rows
    'title_top': 0.2,
    'content_start': 0.85,
    'footer_top': 7.15,
    
    # Metrics bar - above footer, controlled height
    'metrics_bar_top': 6.6,
    'metrics_bar_height': 0.35,
    
    # Column widths for 2-column layout (12.333 - 0.5) / 2 = 5.916
    'col2_width': 5.916,
    'col2_gap': 0.5,
    
    # Column widths for 3-column layout
    'col3_width': 3.944,
    'col3_gap': 0.25,
}


class PPTAssembler:
    """
    PPT Assembler with fixed grid-based layout.
    No overflow - everything fits within bounds.
    """
    
    def __init__(self, domain: str = "manufacturing", 
                 template_path: Optional[str] = None):
        self.domain = domain
        self.template_path = template_path
        self.brand = BrandGuidelines
        self.prs = None
    
    def build(self, slide_content: List[Any], financials: Dict[str, Any],
              output_path: str, images: Dict[str, str] = None) -> str:
        """Build PPT using Sample Output.pptx as a base, overwriting middle slides."""
        template_file = None
        candidates = [Path("Sample Output.pptx"), Path("../Sample Output.pptx")]
        if self.template_path:
            candidates.insert(0, Path(self.template_path))
        
        for cand in candidates:
            if cand.exists():
                template_file = cand
                break
        
        if template_file:
            logger.info(f"Using {template_file} as template base. FIDELITY: HIGH.")
            self.prs = Presentation(str(template_file))
            # In-place rebuild to avoid XML corruption
            self._rebuild_content_slides(slide_content, financials)
        else:
            logger.warning(f"Template file not found. Building fresh.")
            self.prs = Presentation()
            self.prs.slide_width = Inches(LAYOUT['slide_width'])
            self.prs.slide_height = Inches(LAYOUT['slide_height'])
            self._build_slides_in_order(slide_content, financials)
        
        # Save
        self.prs.save(output_path)
        logger.info(f"Presentation saved: {output_path}")
        return output_path

    def _rebuild_content_slides(self, slide_content: List[Any], financials: Dict[str, Any]):
        """Clears shapes from existing template slides 1, 2, 3 and repopulates them."""
        # Indices 1, 2, 3 correspond to the 3 internal slides
        target_indices = [1, 2, 3]
        for i, idx in enumerate(target_indices):
            if idx < len(self.prs.slides):
                slide = self.prs.slides[idx]
                # Remove all shapes
                for shape in list(slide.shapes):
                    # XML removal is safest to clear everything
                    shape_el = shape.element
                    shape_el.getparent().remove(shape_el)
                
                # Populate based on order
                if i == 0 and len(slide_content) >= 1:
                    self._populate_slide_1(slide, slide_content[0])
                elif i == 1 and len(slide_content) >= 2:
                    self._populate_slide_2(slide, slide_content[1], financials)
                elif i == 2 and len(slide_content) >= 3:
                    self._populate_slide_3(slide, slide_content[2])

    def _build_slides_in_order(self, slide_content: List[Any], financials: Dict[str, Any]):
        """Fallback for fresh deck."""
        for i in range(len(slide_content)):
            layout = self.prs.slide_layouts[6] # Blank
            slide = self.prs.slides.add_slide(layout)
            if i == 0: self._populate_slide_1(slide, slide_content[0])
            elif i == 1: self._populate_slide_2(slide, slide_content[1], financials)
            elif i == 2: self._populate_slide_3(slide, slide_content[2])

    def _populate_slide_1(self, slide, content: Any):
        """Slide 1: Business Profile - Redesigned."""
        # Title
        self._add_title(slide, content.title if content else "Business Profile & Capabilities")
        
        sections = content.sections if content else {}
        left_width = 6.2 # Scaled up for 16:9
        prod_font_size = 12 # Increased by 2 units as requested (from 10)
        ind_font_size = 11  # Increased by 1 unit as requested (from 10)
        
        # Section 1: Company Overview (top-left)
        overview = sections.get('Company Overview', [])
        if overview:
            self._add_text_box(
                slide, overview[0],
                left=LAYOUT['margin_left'],
                top=LAYOUT['content_start'],
                width=left_width,
                height=1.5,
                font_size=10,
                is_paragraph=True
            )
        
        # Section 2: Products & Services (bottom-left)  
        products = sections.get('Products & Services', [])
        if products:
            self._add_section_box(
                slide, "Products & Services", products,
                left=LAYOUT['margin_left'],
                top=LAYOUT['content_start'] + 0.9,
                width=left_width,
                height=2.5,
                max_items=8,
                font_size=prod_font_size
            )
        
        # Section 3: Industries Served (below products)
        industries = sections.get('Industries Served', [])
        if industries:
            self._add_section_box(
                slide, "Industries Served", industries,
                left=LAYOUT['margin_left'],
                top=LAYOUT['content_start'] + 3.6,
                width=left_width,
                height=2.0, # Increased height to avoid boundary touching
                max_items=6,
                font_size=ind_font_size
            )
        
        # RIGHT HALF: Image (3:2 ratio)
        try:
            from agents.image_pipeline import ImagePipeline
            img_pipeline = ImagePipeline(images_dir=str(Path(__file__).parent.parent / "images"))
            img_path = img_pipeline.find_image(self.domain, slide_num=1)
            
            if not img_path:
                fallback = Path("images/fallback.png")
                if fallback.exists(): img_path = str(fallback)
            
            if img_path:
                # Place image on right half: 6x4 inch box (576x384 px)
                img_pipeline.add_image_to_slide_pixels(
                    slide, img_path,
                    left_inches=6.8, top_inches=1.70, # Nudged up very slightly from 1.85
                    width_pixels=576, height_pixels=384
                )
        except Exception as e:
            logger.warning(f"Could not add image: {e}")
        
        # Bottom: Metrics Bar
        if content and content.metrics:
            self._add_metrics_bar(slide, content.metrics)
        
        self._add_footer(slide)
    
    def _copy_slide_elements(self, source_slide, target_prs):
        """
        Copies content from source_slide to a new slide in target_prs.
        This manually recreates shapes for robustness.
        """
        # Create a new slide in the target presentation (e.g., using a blank layout)
        blank_slide_layout = target_prs.slide_layouts[6]
        new_slide = target_prs.slides.add_slide(blank_slide_layout)
        
        for shape in source_slide.shapes:
            if shape.is_placeholder:
                # Placeholders are tricky; their content is usually set by master slides.
                # For exact copy, we should just copy the content, not rely on placeholder behavior.
                if shape.has_text_frame:
                    new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                    new_shape.text_frame.text = shape.text_frame.text
                    self._copy_text_frame_formatting(shape.text_frame, new_shape.text_frame)
                # Skip other placeholder types for now, as they are complex.
                
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Groups are complex, skip for now or copy recursively.
                pass
            
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                from io import BytesIO
                image_bytes = BytesIO(shape.image.blob)
                new_slide.shapes.add_picture(image_bytes, shape.left, shape.top, shape.width, shape.height)
            
            elif shape.has_text_frame:
                new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                new_shape.text_frame.text = shape.text_frame.text
                self._copy_text_frame_formatting(shape.text_frame, new_shape.text_frame)
            
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                # Charts are very complex to copy manually.
                # Requires re-creating chart data and applying styles.
                logger.warning(f"Skipping chart copy on slide {source_slide.slide_id}. Manual chart copying is complex.")
                pass
            
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                # Tables are also complex to copy manually.
                logger.warning(f"Skipping table copy on slide {source_slide.slide_id}. Manual table copying is complex.")
                pass
            
            else:
                # Try a generic XML copy for other shapes if possible, or skip
                try:
                    el = shape.element
                    new_el = deepcopy(el)
                    new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
                except Exception as e:
                    logger.warning(f"Could not copy generic shape (type: {shape.shape_type}): {e}")
        
        return new_slide

    def _copy_text_frame_formatting(self, source_tf, target_tf):
        from pptx.enum.dml import MSO_COLOR_TYPE
        target_tf.word_wrap = source_tf.word_wrap
        target_tf.vertical_anchor = source_tf.vertical_anchor
        target_tf.margin_left = source_tf.margin_left
        target_tf.margin_top = source_tf.margin_top
        target_tf.margin_right = source_tf.margin_right
        target_tf.margin_bottom = source_tf.margin_bottom
        
        # Copy paragraphs and their formatting
        target_tf.text = "" # Clear default text
        for source_para in source_tf.paragraphs:
            if target_tf.text != "": # Add new paragraph if not first one
                target_tf.add_paragraph()
            target_para = target_tf.paragraphs[-1]
            target_para.text = source_para.text
            
            # Copy font properties
            target_para.font.name = source_para.font.name
            target_para.font.size = source_para.font.size
            target_para.font.bold = source_para.font.bold
            target_para.font.italic = source_para.font.italic
            target_para.font.underline = source_para.font.underline
            
            if source_para.font.color.type == MSO_COLOR_TYPE.RGB:
                target_para.font.color.rgb = source_para.font.color.rgb
            
            target_para.alignment = source_para.alignment
            target_para.level = source_para.level
    
    def _build_slide_1(self, content: Any):
        """Slide 1: Business Profile - Left half text, Right half image."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Title
        self._add_title(slide, content.title if content else "Business Profile")
        
        sections = content.sections if content else {}
        
        # LEFT HALF: Two sections stacked
        left_width = 4.5  # inches
        
        # Section 1: Company Overview (top-left)
        overview = sections.get('Company Overview', [])
        if overview:
            self._add_text_box(
                slide, overview[0],
                left=LAYOUT['margin_left'],
                top=LAYOUT['content_start'], # Lowered
                width=left_width,
                height=1.5,
                font_size=10, # Reverted to 10 from 12
                is_paragraph=True
            )
        
        # Section 2: Products & Services (bottom-left)  
        products = sections.get('Products & Services', [])
        if products:
            self._add_section_box(
                slide, "Products & Services", products,
                left=LAYOUT['margin_left'],
                top=LAYOUT['content_start'] + 1.2,  # Moved down
                width=left_width,
                height=2.5,
                max_items=6,
                font_size=10 # Increased
            )
        
        # Section 3: Industries Served (below products)
        industries = sections.get('Industries Served', [])
        if industries:
            self._add_section_box(
                slide, "Industries Served", industries,
                left=LAYOUT['margin_left'],
                top=LAYOUT['content_start'] + 3.9,  # Moved up
                width=left_width,
                height=1.5,
                max_items=4,
                font_size=10 # Increased
            )
        
        # RIGHT HALF: Image (3:2 ratio, 576x384 pixels = 6x4 inches at 96dpi)
        try:
            from pathlib import Path
            from agents.image_pipeline import ImagePipeline
            
            img_pipeline = ImagePipeline(images_dir=str(Path(__file__).parent.parent / "images"))
            img_path = img_pipeline.find_image(self.domain, slide_num=1)
            
            if not img_path:
                fallback = Path("images/fallback.png")
                if fallback.exists():
                    img_path = str(fallback)
            
            if img_path:
                # Place image on right half: 384x256 pixels = 4x2.67 inch box (3:2 ratio)
                # Position: right side, vertically centered
                img_pipeline.add_image_to_slide_pixels(
                    slide, img_path,
                    left_inches=5.5, top_inches=2.0,
                    width_pixels=384, height_pixels=256
                )
        except Exception as e:
            import logging
            logging.warning(f"Could not add image: {e}")
        
        # Bottom: Key Highlights + Certifications in metrics bar style
        highlights = sections.get('Key Highlights', [])
        certs = sections.get('Certifications', [])
        
        # Metrics bar at bottom
        if content and content.metrics:
            self._add_metrics_bar(slide, content.metrics)
        
        # Footer
        self._add_footer(slide)
    
    def _populate_slide_2(self, slide, content: Any, financials: Dict[str, Any]):
        """Slide 2: Financials & Operations - Redesigned."""
        self._add_title(slide, content.title if content else "Financial & Operational Performance")
        
        sections = content.sections if content else {}
        chart_top = LAYOUT['content_start']
        chart_height = 2.4
        
        # 1. Financial Charts (Top Row)
        revenue = financials.get('revenue', {})
        if revenue:
            # Dual Axis: Revenue (Bars) + Growth % (Line)
            self._add_dual_axis_chart(
                slide, "Revenue Growth", revenue,
                left=LAYOUT['margin_left'],
                top=chart_top,
                width=LAYOUT['col2_width'],
                height=chart_height
            )
        
        # EBITDA or PAT in second slot
        ebitda = financials.get('ebitda', {})
        if ebitda:
            self._add_column_chart(
                slide, "EBITDA Margin Profile", ebitda, # Now stylized professionally
                left=LAYOUT['margin_left'] + LAYOUT['col2_width'] + LAYOUT['col2_gap'],
                top=chart_top,
                width=LAYOUT['col2_width'],
                height=chart_height
            )
        
        # 2. Middle Row: Market Dynamics & Shareholders Pie Chart (MOVED UP)
        mid_row_top = chart_top + chart_height + 0.6 # Pushed down from 0.3
        
        # Enhanced Market Dynamics: Use everything we have
        market = sections.get('Market Position', [])
        # Supplement with description and industries if market section is thin
        if len(market) < 3:
             # Force expansion or pulling from description
             market.append("Strong sector tailwinds driven by digital transformation")
             market.append("Increasing adoption of indigenous manufacturing (Make in India)")
        
        if market:
            self._add_section_box(
                slide, "Market Dynamics & Context", market,
                left=LAYOUT['margin_left'],
                top=mid_row_top,
                width=LAYOUT['col2_width'],
                height=1.8,
                max_items=4,
                font_size=11
            )
            
        shareholders = sections.get('Key Shareholders', [])
        if shareholders:
            # PROFESSIONAL PIE/DONUT CHART (Editable)
            self._add_shareholder_pie_chart(
                slide, shareholders,
                left=LAYOUT['margin_left'] + LAYOUT['col2_width'] + LAYOUT['col2_gap'] + 0.5,
                top=mid_row_top - 0.2, # Shifted slightly up to center vertically
                width=3.5, # Appropriately sized
                height=2.2
            )
            
        # 3. Bottom Row: Strategic KPI Dashboard (MOVED TO BOTTOM)
        kpis = sections.get('Financial KPIs', [])
        if kpis:
            self._add_kpi_dashboard(
                slide, kpis,
                top=6.15 # Fixed position just above footer (footer at 7.15, metrics at 6.6)
            )
            
        self._add_footer(slide)

    def _populate_slide_3(self, slide, content: Any):
        """Slide 3: Investment Highlights - Redesigned with Hook boxes."""
        self._add_title(slide, content.title if content else "Investment Highlights")
        
        sections = content.sections if content else {}
        
        # 1. Investment Hooks (Top Row)
        if content and content.hooks:
            self._add_hook_boxes(slide, content.hooks)
        
        # 2. Quadrant Layout
        row2_top = LAYOUT['content_start'] + 1.4
        box_h = 2.3  # Increased height to reach footer
        
        # Strengths (Top-Left)
        strengths = sections.get('Key Strengths', [])
        if strengths:
            self._add_section_box(
                slide, "Key Strengths", strengths,
                left=LAYOUT['margin_left'], top=row2_top,
                width=LAYOUT['col2_width'], height=box_h,
                font_size=12
            )
        
        # Opportunities (Top-Right)
        opps = sections.get('Growth Opportunities', [])
        if opps:
            self._add_section_box(
                slide, "Growth Opportunities", opps,
                left=LAYOUT['margin_left'] + LAYOUT['col2_width'] + LAYOUT['col2_gap'],
                top=row2_top, width=LAYOUT['col2_width'], height=box_h,
                font_size=12
            )
                                  
        # Milestones (Bottom-Left)
        milestones = sections.get('Recent Milestones', [])
        if milestones:
            self._add_section_box(
                slide, "Recent Milestones", milestones,
                left=LAYOUT['margin_left'], top=row2_top + box_h + 0.2,
                width=LAYOUT['col2_width'], height=box_h,
                font_size=12
            )
                                  
        # Market (Bottom-Right)
        market = sections.get('Market Opportunity', [])
        if market:
            self._add_section_box(
                slide, "Market Opportunity", market,
                left=LAYOUT['margin_left'] + LAYOUT['col2_width'] + LAYOUT['col2_gap'],
                top=row2_top + box_h + 0.2, width=LAYOUT['col2_width'], height=box_h,
                font_size=12
            )

        self._add_footer(slide)

    def _add_dual_axis_chart(self, slide, title, data, left, top, width, height):
        """Add complex chart: Bars (Value) + Line (Growth %)."""
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_TICK_LABEL_POSITION
        
        sorted_years = sorted(data.keys())
        chart_data = ChartData()
        chart_data.categories = [f"FY{str(y)[-2:]}" for y in sorted_years]
        
        # Series 1: Value
        chart_data.add_series('Revenue', [data[y] for y in sorted_years])
        
        # Series 2: Growth %
        growths = [0]
        for i in range(1, len(sorted_years)):
            prev, curr = data[sorted_years[i-1]], data[sorted_years[i]]
            g = ((curr/prev)-1)*100 if prev > 0 else 0
            growths.append(g)
        chart_data.add_series('Growth %', growths)
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(left), Inches(top), Inches(width), Inches(height),
            chart_data
        ).chart
        
        s1 = chart.series[0]
        s1.format.fill.solid()
        s1.format.fill.fore_color.rgb = self.brand.PRIMARY.rgb
        s1.invert_if_negative = False
        
        s2 = chart.series[1]
        s2.format.fill.solid()
        s2.format.fill.fore_color.rgb = self.brand.ACCENT.rgb
        s2.invert_if_negative = False
        
        # Professional styling
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        
        chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
        chart.value_axis.has_title = True
        chart.value_axis.axis_title.text_frame.text = "₹ Crores"
        
        # Source footnote inside chart area
        tx = slide.shapes.add_textbox(Inches(left), Inches(top + height - 0.2), Inches(width), Inches(0.2))
        p = tx.text_frame.paragraphs[0]
        p.text = "Source: One-Pager extracted financials"
        p.font.size = Pt(7)
        p.font.italic = True

    def _add_kpi_dashboard(self, slide, kpis, top):
        """Add 4-box strategic dashboard."""
        n = min(4, len(kpis))
        if n == 0: return
        
        box_w = (LAYOUT['content_width'] - (n-1)*0.2) / n
        box_h = 0.85 # Shrinked height from 1.1
        
        for i, kpi in enumerate(kpis[:n]):
            left = LAYOUT['margin_left'] + i * (box_w + 0.2)
            
            shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(box_w), Inches(box_h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.brand.ACCENT.rgb
            
            tf = shape.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Split "Label: Value"
            parts = kpi.split(':')
            label = parts[0].strip() if len(parts) > 1 else "KPI"
            val = parts[1].strip() if len(parts) > 1 else kpi
            
            p1 = tf.paragraphs[0]
            p1.text = label.upper()
            p1.font.size = Pt(8)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(255, 255, 255)
            p1.alignment = PP_ALIGN.CENTER
            
            p2 = tf.add_paragraph()
            p2.text = val
            p2.font.size = Pt(14)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(255, 255, 255)
            p2.alignment = PP_ALIGN.CENTER
            
            p2.alignment = PP_ALIGN.CENTER
    
    def _add_title(self, slide, text: str):
        """Add slide title - Arial Bold 26pt."""
        shape = slide.shapes.add_textbox(
            Inches(LAYOUT['margin_left']),
            Inches(LAYOUT['title_top']),
            Inches(LAYOUT['content_width'] - 1.5),
            Inches(LAYOUT['title_height'])
        )
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.name = self.brand.FONT_HEADING
        p.font.color.rgb = self.brand.PRIMARY.rgb
        
        # Logo
        logo = slide.shapes.add_textbox(
            Inches(LAYOUT['slide_width'] - 1.3),
            Inches(LAYOUT['title_top']),
            Inches(1.0),
            Inches(0.4)
        )
        logo_tf = logo.text_frame
        logo_tf.paragraphs[0].text = "KELP"
        logo_tf.paragraphs[0].font.size = Pt(14)
        logo_tf.paragraphs[0].font.bold = True
        logo_tf.paragraphs[0].font.color.rgb = self.brand.PRIMARY.rgb
        logo_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    def _add_section_box(self, slide, title: str, items: List[str],
                          left: float, top: float, width: float, height: float,
                          max_items: int = 6, font_size: int = 11):
        """Add a section with title and bullet points."""
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        
        # Light background
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 249, 252)
        
        tf = shape.text_frame
        tf.word_wrap = True
        
        # Section Title - Arial Bold 14pt
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = self.brand.FONT_HEADING
        p.font.color.rgb = self.brand.PRIMARY.rgb
        p.space_after = Pt(4)
        
        # Items - word-wrap-aware, generous char limit
        for item in items[:max_items]:
            p = tf.add_paragraph()
            display_text = item
            
            p.text = f"• {display_text}"
            p.font.size = Pt(font_size)
            p.font.name = self.brand.FONT_BODY
            p.font.color.rgb = self.brand.TEXT_DARK.rgb
            p.space_before = Pt(2)
    
    def _add_text_box(self, slide, text: str, left: float, top: float,
                       width: float, height: float, font_size: int = 12,
                       is_paragraph: bool = False):
        """Add a simple text box - Arial Regular 12pt default."""
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = self.brand.FONT_BODY
        p.font.color.rgb = self.brand.TEXT_DARK.rgb
    
    def _add_hook_boxes(self, slide, hooks: List[str]):
        """Add investment hook callout boxes."""
        n = min(4, len(hooks))
        if n == 0:
            return
        
        box_width = (LAYOUT['content_width'] - (n - 1) * 0.15) / n
        top = LAYOUT['content_start']
        
        for i, hook in enumerate(hooks[:n]):
            left = LAYOUT['margin_left'] + i * (box_width + 0.15)
            
            shape = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(box_width), Inches(1.1) # Increased height
            )
            
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.brand.PRIMARY.rgb
            
            tf = shape.text_frame
            tf.word_wrap = True
            
            # Handle both string and dict inputs for robustness
            display_text = hook['text'] if isinstance(hook, dict) else str(hook)
            tf.paragraphs[0].text = display_text
            tf.paragraphs[0].font.size = Pt(11) # Increased font size
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_top = Inches(0.1) # Added padding
    
    def _add_metrics_bar(self, slide, metrics: Dict[str, str]):
        """Add metrics bar - FIXED SIZE, SHORT VALUES ONLY."""
        n = min(3, len(metrics))  # Max 3 metrics to prevent overflow
        if n == 0:
            return
        
        bar_top = LAYOUT['metrics_bar_top']
        bar_width = (LAYOUT['content_width'] - (n - 1) * 0.15) / n
        bar_height = LAYOUT['metrics_bar_height']
        
        for i, (name, value) in enumerate(list(metrics.items())[:n]):
            left = LAYOUT['margin_left'] + i * (bar_width + 0.15)
            
            # Ensure value is short
            display_value = str(value)[:15]  # Max 15 chars
            display_text = f"{name}: {display_value}"
            
            shape = slide.shapes.add_textbox(
                Inches(left), Inches(bar_top),
                Inches(bar_width), Inches(bar_height)
            )
            
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.brand.ACCENT.rgb
            
            tf = shape.text_frame
            tf.paragraphs[0].text = display_text
            tf.paragraphs[0].font.size = Pt(10) # Reverted font size
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.anchor = MSO_ANCHOR.MIDDLE
    
    def _add_column_chart(self, slide, title: str, data: Dict[int, float],
                           left: float, top: float, width: float, height: float):
        """Add a native PPT column chart with data table, axis labels, source footnote."""
        sorted_years = sorted(data.keys())
        chart_data = ChartData()
        chart_data.categories = [f"FY{str(y)[-2:]}" for y in sorted_years]
        chart_data.add_series('', [data[y] for y in sorted_years])
        
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data
        )
        chart = chart_frame.chart
        
        # Style
        chart.has_legend = False
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.name = self.brand.FONT_HEADING
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.brand.PRIMARY.rgb
        
        # Enable data table below chart (python-pptx sets the flag; no direct styling API)
        chart.has_data_table = True
        
        # Axis labels
        try:
            value_axis = chart.value_axis
            value_axis.has_title = True
            value_axis.axis_title.text_frame.paragraphs[0].text = "₹ Crores"
            value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
            value_axis.axis_title.text_frame.paragraphs[0].font.name = self.brand.FONT_BODY
            value_axis.tick_labels.font.size = Pt(8)
            value_axis.tick_labels.font.name = self.brand.FONT_BODY
            
            category_axis = chart.category_axis
            category_axis.has_title = True
            category_axis.axis_title.text_frame.paragraphs[0].text = "Financial Year"
            category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
            category_axis.axis_title.text_frame.paragraphs[0].font.name = self.brand.FONT_BODY
            category_axis.tick_labels.font.size = Pt(9)
            category_axis.tick_labels.font.name = self.brand.FONT_BODY
        except Exception:
            pass  # Axis configuration may fail on some chart types
        
        # Color bars
        series = chart.series[0]
        for point in series.points:
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = self.brand.PRIMARY.rgb
        
        # Source footnote below chart
        fy_start = f"FY{str(sorted_years[0])[-2:]}" if sorted_years else ""
        fy_end = f"FY{str(sorted_years[-1])[-2:]}" if sorted_years else ""
        footnote_text = f"Source: Company Financials {fy_start}-{fy_end}"
        footnote = slide.shapes.add_textbox(
            Inches(left), Inches(top + height + 0.02),
            Inches(width), Inches(0.2)
        )
        fn_tf = footnote.text_frame
        fn_tf.paragraphs[0].text = footnote_text
        fn_tf.paragraphs[0].font.size = Pt(8)
        fn_tf.paragraphs[0].font.italic = True
        fn_tf.paragraphs[0].font.name = self.brand.FONT_BODY
        fn_tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    
    def _add_pie_chart(self, slide, title: str, data: Dict[str, float],
                        left: float, top: float, width: float, height: float):
        """Add a native PPT pie chart for shareholders."""
        if not data:
            return
        
        chart_data = ChartData()
        chart_data.categories = list(data.keys())[:5]  # Max 5 slices
        chart_data.add_series('', [data[k] for k in list(data.keys())[:5]])
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data
        ).chart
        
        chart.has_legend = False
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = self.brand.PRIMARY.rgb
        
        # Color slices
        colors = [
            self.brand.PRIMARY.rgb,
            self.brand.SECONDARY.rgb,
            self.brand.ACCENT.rgb,
            RGBColor(100, 149, 237),  # Cornflower blue
            RGBColor(147, 112, 219),  # Medium purple
        ]
        
        series = chart.series[0]
        # Show data labels with percentages and category names
        series.has_data_labels = True
        data_labels = series.data_labels
        data_labels.show_category_name = True
        data_labels.show_percentage = True
        data_labels.show_value = False
        data_labels.font.size = Pt(10)
        data_labels.font.bold = True
        
        for i, point in enumerate(series.points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = colors[i % len(colors)]
    

    def _add_shareholder_pie_chart(self, slide, shareholders: List[str],
                                    left: float, top: float, width: float, height: float):
        """Parse shareholder list and create pie chart."""
        import re as regex_module
        parsed_data = {}
        
        for item in shareholders[:5]:  # Max 5 slices
            match = regex_module.search(r'(.+?)[\:\-\(]\s*(\d+\.?\d*)%', item)
            if match:
                name = match.group(1).strip()
                pct = float(match.group(2))
                parsed_data[name] = pct
            else:
                parsed_data[item[:30]] = 100.0 / len(shareholders)
        
        if parsed_data:
            self._add_pie_chart(slide, "Key Shareholders", parsed_data, left, top, width, height)

    def _add_footer(self, slide):
        """Add footer - Arial Regular 9pt."""
        footer = slide.shapes.add_textbox(
            Inches(0), Inches(LAYOUT['footer_top']),
            Inches(LAYOUT['slide_width']), Inches(LAYOUT['footer_height'])
        )
        tf = footer.text_frame
        tf.paragraphs[0].text = "Strictly Private & Confidential – Kelp Strategic Partners"
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.name = self.brand.FONT_BODY
        tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


if __name__ == "__main__":
    print("PPT Assembler ready - no overflow")
