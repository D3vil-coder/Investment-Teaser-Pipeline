"""
Image Pipeline - Intelligent image selection and placement
"""

from pathlib import Path
from PIL import Image
from pptx.util import Inches
import logging

logger = logging.getLogger(__name__)


class ImagePipeline:
    """
    Handles intelligent image selection, resizing, and placement in PPTs.
    """
    
    def __init__(self, images_dir: str = "images"):
        self.images_dir = Path(images_dir)
        self.image_cache = {}
        self._scan_images()
    
    def _scan_images(self):
        """Scan images directory and catalog by filename."""
        if not self.images_dir.exists():
            logger.warning(f"Images directory not found: {self.images_dir}")
            return
        
        # Scan root level images
        for img_file in self.images_dir.glob('*'):
            if img_file.is_file() and img_file.suffix.lower() in ['.jpg', '.jpeg', '.png']:
                key = img_file.stem.lower().replace('_', '').replace('-', '')
                self.image_cache[key] = str(img_file)
        
        # Scan subdirectories (domain folders like technology/, logistics/)
        for subdir in self.images_dir.iterdir():
            if subdir.is_dir():
                domain_name = subdir.name.lower()
                for img_file in subdir.glob('*'):
                    if img_file.suffix.lower() in ['.jpg', '.jpeg', '.png']:
                        # Store with domain prefix: "technology_officemoddern"
                        key = f"{domain_name}_{img_file.stem.lower().replace('_', '').replace('-', '')}"
                        self.image_cache[key] = str(img_file)
                        # Also store domain as a general key (first image in folder)
                        if domain_name not in self.image_cache:
                            self.image_cache[domain_name] = str(img_file)
        
        logger.info(f"Found {len(self.image_cache)} images in library")
    
    def find_image(self, domain: str, slide_num: int = 1) -> str:
        """
        Find best matching image for domain and slide.
        
        Args:
            domain: Domain name (e.g., 'technology', 'Technology & IT Services')
            slide_num: Slide number (1-3)
        
        Returns:
            Path to image file, or None if not found
        """
        # Normalize domain name
        domain_lower = domain.lower()
        
        # Map full domain names to folder names
        domain_map = {
            'technology & it services': 'technology',
            'technology': 'technology',
            'it services': 'technology',
            'manufacturing': 'manufacturing',
            'logistics & supply chain': 'logistics',
            'logistics': 'logistics',
            'supply chain': 'logistics',
            'consumer & retail': 'consumer',
            'consumer': 'consumer',
            'retail': 'consumer',
            'healthcare & pharma': 'healthcare',
            'healthcare': 'healthcare',
            'pharma': 'healthcare',
            'infrastructure': 'infrastructure',
            'chemicals': 'chemicals',
            'automotive': 'automotive',
        }
        
        # Get normalized domain folder name
        folder_name = domain_map.get(domain_lower, domain_lower.split()[0].lower())
        
        # Check if domain folder exists and has images
        domain_folder = self.images_dir / folder_name
        if domain_folder.exists() and domain_folder.is_dir():
            images = list(domain_folder.glob('*.png')) + list(domain_folder.glob('*.jpg')) + list(domain_folder.glob('*.jpeg'))
            if images:
                import random
                # Return random image from domain folder
                selected = str(random.choice(images))
                logger.info(f"Found domain image: {selected}")
                return selected
        
        # Fallback: search cache
        search_terms = [
            folder_name,
            f"{folder_name}{slide_num}",
            f"{folder_name}main",
        ]
        
        for term in search_terms:
            normalized = term.lower().replace('_', '').replace('-', '')
            if normalized in self.image_cache:
                logger.info(f"Found cached image for {domain}: {self.image_cache[normalized]}")
                return self.image_cache[normalized]
        
        logger.warning(f"No image found for domain={domain} (folder={folder_name})")
        return None
    
    def add_image_to_slide(self, slide, image_path: str, 
                          left: float, top: float, 
                          max_width: float, max_height: float):
        """
        Add image to slide with intelligent resizing (maintains 3:2 aspect ratio).
        
        Args:
            slide: PPT slide object
            image_path: Path to image file
            left, top: Position in inches
            max_width, max_height: Maximum dimensions in inches
        """
        if not image_path or not Path(image_path).exists():
            logger.warning(f"Image not found: {image_path}")
            return False
        
        try:
            # Open image to get dimensions
            with Image.open(image_path) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height
            
            # Calculate scaled dimensions (maintain aspect ratio)
            if aspect_ratio > (max_width / max_height):
                # Width-constrained
                final_width = max_width
                final_height = max_width / aspect_ratio
            else:
                # Height-constrained
                final_height = max_height
                final_width = max_height * aspect_ratio
            
            # Add image centered in max bounds
            left_centered = left + (max_width - final_width) / 2
            top_centered = top + (max_height - final_height) / 2
            
            slide.shapes.add_picture(
                image_path,
                Inches(left_centered),
                Inches(top_centered),
                width=Inches(final_width),
                height=Inches(final_height)
            )
            
            logger.info(f"Added image {Path(image_path).name} at ({left:.2f}, {top:.2f}) "
                       f"size {final_width:.2f}x{final_height:.2f}")
            return True
            
        except Exception as e:
            logger.error(f"Error adding image {image_path}: {e}")
            return False

    def add_image_to_slide_pixels(self, slide, image_path: str,
                                   left_inches: float, top_inches: float,
                                   width_pixels: int, height_pixels: int):
        """
        Add image to slide scaled to EXACT pixel dimensions (e.g., 1200x800).
        Maintains aspect ratio by fitting within the box.
        
        Args:
            slide: PPT slide object
            image_path: Path to image file
            left_inches, top_inches: Position in inches
            width_pixels, height_pixels: Target box size in pixels
        """
        if not image_path or not Path(image_path).exists():
            logger.warning(f"Image not found: {image_path}")
            return False
        
        try:
            # Convert pixels to inches (96 DPI)
            max_width_inches = width_pixels / 96.0
            max_height_inches = height_pixels / 96.0
            
            # Open image to get dimensions
            with Image.open(image_path) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height
            
            # Calculate scaled dimensions to fit in box
            target_ratio = width_pixels / height_pixels
            
            if aspect_ratio > target_ratio:
                # Width-constrained
                final_width = max_width_inches
                final_height = max_width_inches / aspect_ratio
            else:
                # Height-constrained
                final_height = max_height_inches
                final_width = max_height_inches * aspect_ratio
            
            # Center in box
            left_centered = left_inches + (max_width_inches - final_width) / 2
            top_centered = top_inches + (max_height_inches - final_height) / 2
            
            slide.shapes.add_picture(
                image_path,
                Inches(left_centered),
                Inches(top_centered),
                width=Inches(final_width),
                height=Inches(final_height)
            )
            
            logger.info(f"Added image {Path(image_path).name} "
                       f"scaled to fit {width_pixels}x{height_pixels}px box")
            return True
            
        except Exception as e:
            logger.error(f"Error adding image: {e}")
            return False


if __name__ == "__main__":
    # Test
    pipeline = ImagePipeline()
    print(f"Image cache: {len(pipeline.image_cache)} images")
    
    test_img = pipeline.find_image('technology', 1)
    if test_img:
        print(f"Found: {test_img}")
    else:
        print("No image found for technology/1")
